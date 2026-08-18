"""Build a plain WINGS3 deck: default Office layouts, no branded template.

Each slide is title + optional bullets + speaker notes.
Run: python3 scripts/build_wings3_deck.py
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

_SCRIPTS_DIR = Path(__file__).resolve().parent
WINGS3_ROOT = _SCRIPTS_DIR.parent

if str(_SCRIPTS_DIR) not in sys.path:
    sys.path.insert(0, str(_SCRIPTS_DIR))

from wings3_content import EXPECTED_SLIDE_COUNT, SLIDES  # noqa: E402

OUTPUT = WINGS3_ROOT / "MLflow-on-RHOAI-Deep-Dive.pptx"


def _set_notes(slide, text: str) -> None:
    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    notes.text = text.strip()


def _set_title(slide, title: str) -> None:
    slide.shapes.title.text = title


def _add_body(slide, lines: list[str]) -> None:
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0


def _add_title_slide(prs: Presentation, spec: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    _set_title(slide, spec["title"])
    subtitle = spec.get("subtitle", "")
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle
    _set_notes(slide, spec["notes"])
    return slide


def _add_section_slide(prs: Presentation, spec: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _set_title(slide, spec["title"])
    subtitle = spec.get("subtitle")
    if subtitle:
        box = slide.shapes.add_textbox(Inches(0.75), Inches(3.2), Inches(8.5), Inches(1))
        p = box.text_frame.paragraphs[0]
        p.text = subtitle
    _set_notes(slide, spec["notes"])
    return slide


def _add_content_slide(prs: Presentation, spec: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_title(slide, spec["title"])
    _add_body(slide, spec.get("bullets") or [])
    _set_notes(slide, spec["notes"])
    return slide


FILLERS = {
    "title": _add_title_slide,
    "section": _add_section_slide,
    "content": _add_content_slide,
}


def build_deck() -> Path:
    prs = Presentation()
    for spec in SLIDES:
        FILLERS[spec["layout"]](prs, spec)
    if OUTPUT.exists():
        OUTPUT.unlink()
    prs.save(str(OUTPUT))
    verify = Presentation(str(OUTPUT))
    assert len(verify.slides) == EXPECTED_SLIDE_COUNT, (
        f"Expected {EXPECTED_SLIDE_COUNT} slides, got {len(verify.slides)}"
    )
    for i, slide in enumerate(verify.slides, 1):
        assert slide.has_notes_slide, f"Slide {i} missing notes"
        assert slide.notes_slide.notes_text_frame.text.strip(), f"Slide {i} has empty notes"
    return OUTPUT


def main() -> None:
    path = build_deck()
    print(f"Wrote {path}")
    print(f"Slides: {EXPECTED_SLIDE_COUNT}")


if __name__ == "__main__":
    main()
