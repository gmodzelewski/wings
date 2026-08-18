"""Integration test for the plain WINGS3 deck builder."""

from __future__ import annotations

import sys
from pathlib import Path

import pytest
from pptx import Presentation

REPO_ROOT = Path(__file__).resolve().parent.parent
WINGS3_SCRIPTS = REPO_ROOT / "scripts"
sys.path.insert(0, str(WINGS3_SCRIPTS))
sys.path.insert(0, str(REPO_ROOT))

from build_wings3_deck import OUTPUT, build_deck  # noqa: E402
from wings3_content import EXPECTED_SLIDE_COUNT, SLIDES  # noqa: E402

pytest.importorskip("pptx")


def test_build_creates_expected_slides_with_notes():
    out = build_deck()
    assert out == OUTPUT
    assert out.is_file()

    prs = Presentation(str(out))
    assert len(prs.slides) == EXPECTED_SLIDE_COUNT
    assert len(prs.slides) == len(SLIDES)

    for i, slide in enumerate(prs.slides, 1):
        assert slide.has_notes_slide
        notes = slide.notes_slide.notes_text_frame.text.strip()
        assert notes, f"Slide {i} has empty notes"


def test_title_slide_has_headline():
    build_deck()
    prs = Presentation(str(OUTPUT))
    first = prs.slides[0]
    texts = [
        shape.text_frame.text.strip()
        for shape in first.shapes
        if shape.has_text_frame and shape.text_frame.text.strip()
    ]
    assert any("MLflow" in t for t in texts)


def test_every_spec_has_notes_and_title():
    for spec in SLIDES:
        assert spec["title"].strip()
        assert spec["notes"].strip()
        assert spec["layout"] in ("title", "section", "content")
